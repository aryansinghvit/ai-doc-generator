from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import simpleSplit
from fastapi import FastAPI, Depends, HTTPException, status
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from fastapi.security import OAuth2PasswordBearer # <--- NEW
from pydantic import BaseModel
from sqlalchemy.orm import Session
from docx import Document
from pptx import Presentation
from io import BytesIO
from jose import jwt  # <--- NEW

# Custom Imports
import models
import security
import ai_service
from database import engine, get_db

models.Base.metadata.create_all(bind=engine)

app = FastAPI()

oauth2_scheme = OAuth2PasswordBearer(tokenUrl="login")

def get_current_user(token: str = Depends(oauth2_scheme), db: Session = Depends(get_db)):
    credentials_exception = HTTPException(
        status_code=status.HTTP_401_UNAUTHORIZED,
        detail="Could not validate credentials",
        headers={"WWW-Authenticate": "Bearer"},
    )
    try:
        payload = jwt.decode(token, security.SECRET_KEY, algorithms=[security.ALGORITHM])
        email: str = payload.get("sub")
        if email is None:
            raise credentials_exception
    except:
        raise credentials_exception

    user = db.query(models.User).filter(models.User.email == email).first()
    if user is None:
        raise credentials_exception
    return user

# --- NEW SECTION: ALLOW FRONTEND CONNECTION ---
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], # This is where React lives
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
# ----------------------------------------------
# --- DATA SHAPES (Pydantic Models) ---
class UserCreate(BaseModel):
    email: str
    password: str

class UserLogin(BaseModel):
    email: str
    password: str

class Token(BaseModel):
    access_token: str
    token_type: str

class FeedbackRequest(BaseModel):
    is_liked: bool
    is_disliked: bool
    user_notes: str

class RefineRequest(BaseModel):
    content: str
    instruction: str # e.g., "Make it professional"


class GenerateRequest(BaseModel):
    topic: str

@app.get("/")
def read_root():
    return {"message": "AI Document Platform is Running!"}

@app.post("/generate")
def generate_text(request: GenerateRequest):
    # We changed the prompt to ask for a "detailed comprehensive guide"
    prompt = (
        f"Write a detailed and comprehensive document about: '{request.topic}'. "
        f"Include the history, key features, specifications, and impact. "
        f"The content should be approximately 300 words long and professionally formatted."
    )
    
    generated_text = ai_service.generate_document_content(prompt)
    return {"content": generated_text}



# For Word document export


# Define the data structure for the export request
class ExportRequest(BaseModel):
    content: str



@app.post("/export/docx")
def export_docx(request: ExportRequest):
    # 1. Create a new Word Document in memory
    doc = Document()
    doc.add_heading('AI Generated Document', 0)
    doc.add_paragraph(request.content)

    # 2. Save it to a "buffer" (temporary memory) instead of hard disk
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0) # Go back to the start of the file

    # 3. Send the file back to the browser
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": "attachment; filename=generated_doc.docx"}
    )

class ProjectSaveRequest(BaseModel):
    title: str
    content: str
    doc_type: str

@app.post("/export/pdf")
def export_pdf(request: ExportRequest):
    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter

    # 1. Add Title
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, height - 50, "AI Generated Document")

    # 2. Add Content (with text wrapping)
    c.setFont("Helvetica", 12)
    text_object = c.beginText(72, height - 80)

    # Split long paragraphs into lines that fit the page width
    # 460 is roughly the width available between margins
    wrapped_lines = simpleSplit(request.content, "Helvetica", 12, 460)

    for line in wrapped_lines:
        text_object.textLine(line)

        # If we run out of space on the page, make a new page
        if text_object.getY() < 50:
            c.drawText(text_object)
            c.showPage() # New Page
            text_object = c.beginText(72, height - 50)
            text_object.setFont("Helvetica", 12)

    c.drawText(text_object)
    c.save()
    buffer.seek(0)

    return StreamingResponse(
        buffer,
        media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=generated_doc.pdf"}
    )

# For PowerPoint export

@app.post("/export/pptx")
def export_pptx(request: ExportRequest):
    # 1. Create a PowerPoint Presentation
    prs = Presentation()
    
    # 2. Add a Slide (Layout 1 is usually "Title and Content")
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    
    # 3. Put text into the slide
    # We set a generic title for now (or you could send it from frontend)
    if slide.shapes.title:
        slide.shapes.title.text = "AI Generated Slide"
        
    # Put the AI content in the body box
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = request.content
    
    # 4. Save to memory buffer
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    # 5. Send file back
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=generated_pres.pptx"}
    )

# --- AUTHENTICATION ENDPOINTS ---

@app.post("/register", status_code=status.HTTP_201_CREATED)
def register_user(user: UserCreate, db: Session = Depends(get_db)):
    # 1. Check if email already exists
    existing_user = db.query(models.User).filter(models.User.email == user.email).first()
    if existing_user:
        raise HTTPException(status_code=400, detail="Email already registered")
    
    # 2. Hash the password (Security!)
    hashed_pwd = security.get_password_hash(user.password)
    
    # 3. Create new User in Database
    new_user = models.User(email=user.email, password=hashed_pwd)
    db.add(new_user)
    db.commit()
    db.refresh(new_user)
    
    return {"message": "User created successfully", "id": new_user.id}

@app.post("/login", response_model=Token)
def login_for_access_token(user: UserLogin, db: Session = Depends(get_db)):
    # 1. Find user by email
    db_user = db.query(models.User).filter(models.User.email == user.email).first()
    if not db_user:
        raise HTTPException(status_code=400, detail="Incorrect email or password")
    
    # 2. Check if password matches
    if not security.verify_password(user.password, db_user.password):
        raise HTTPException(status_code=400, detail="Incorrect email or password")
    
    # 3. Generate the "ID Card" (Token)
    access_token = security.create_access_token(data={"sub": db_user.email})
    
    return {"access_token": access_token, "token_type": "bearer"}

@app.post("/projects")
def save_project(
    project: ProjectSaveRequest, 
    db: Session = Depends(get_db), 
    current_user: models.User = Depends(get_current_user)
):
    # 1. Create the Project Entry
    new_proj = models.Project(
        title=project.title, 
        doc_type=project.doc_type, 
        user_id=current_user.id
    )
    db.add(new_proj)
    db.commit()
    db.refresh(new_proj)

    # 2. Save the Content as a Section
    # (For simplicity, we save the whole text as Section 1)
    new_section = models.DocumentSection(
        project_id=new_proj.id,
        sequence_order=1,
        heading="Main Content",
        content=project.content
    )
    db.add(new_section)
    db.commit()

    return {"message": "Project Saved Successfully!", "project_id": new_proj.id}


# --- HISTORY & DELETE ENDPOINTS ---

@app.get("/projects")
def get_user_projects(db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    # 1. Get all projects for this user, ordered by newest first
    projects = db.query(models.Project).filter(models.Project.user_id == current_user.id).order_by(models.Project.created_at.desc()).all()
    
    # 2. Format the data to send back (including the content!)
    results = []
    for p in projects:
        # Get the main content (first section)
        content_section = db.query(models.DocumentSection).filter(models.DocumentSection.project_id == p.id).first()
        results.append({
            "id": p.id,
            "title": p.title,
            "created_at": p.created_at,
            "content": content_section.content if content_section else ""
        })
    return results

@app.delete("/projects/{project_id}")
def delete_project(project_id: int, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    # 1. Find the project (and make sure it belongs to this user!)
    project = db.query(models.Project).filter(models.Project.id == project_id, models.Project.user_id == current_user.id).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # 2. Delete the sections first (to keep database clean)
    db.query(models.DocumentSection).filter(models.DocumentSection.project_id == project_id).delete()
    
    # 3. Delete the project
    db.delete(project)
    db.commit()
    return {"message": "Project deleted"}

# --- FEEDBACK & REFINEMENT ENDPOINTS ---

@app.put("/projects/{project_id}/feedback")
def update_feedback(project_id: int, feedback: FeedbackRequest, db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    # Find the main section of the project
    section = db.query(models.DocumentSection).join(models.Project).filter(
        models.Project.id == project_id,
        models.Project.user_id == current_user.id
    ).first()
    
    if not section:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Update fields
    section.is_liked = feedback.is_liked
    section.is_disliked = feedback.is_disliked
    section.user_notes = feedback.user_notes
    db.commit()
    return {"message": "Feedback saved"}

@app.post("/refine")
def refine_text(request: RefineRequest):
    # Ask AI to rewrite the text based on instruction
    prompt = (
        f"Original Text: '{request.content}'\n\n"
        f"Instruction: {request.instruction}\n\n"
        f"Rewrite the text following the instruction. Keep it detailed."
    )
    refined_text = ai_service.generate_document_content(prompt)
    return {"content": refined_text}